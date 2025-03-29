from abc import ABC, abstractmethod
from typing import Dict, List, Any, Optional, Union
import os
import json
import time
import logging
import urllib.parse
import random
from datetime import datetime
import re
import shutil

# Configure logging - only log to file by default, not to console
file_handler = logging.FileHandler("agent_debug.log")
file_handler.setLevel(logging.INFO)
file_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))

# Set up the root logger with just the file handler
logging.basicConfig(
    level=logging.INFO,
    handlers=[file_handler]
)
logger = logging.getLogger("tools")

# Try to import PyPDF2 for PDF file processing
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    logger.warning("PyPDF2 not installed. Will not be able to read PDF files.")
    PYPDF2_AVAILABLE = False

# Try to import docx2txt for DOCX file processing
try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except ImportError:
    logger.warning("docx2txt not installed. Will not be able to read DOCX files.")
    DOCX2TXT_AVAILABLE = False

# Try to import python-pptx for PPTX file processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not installed. Will not be able to read PPTX files.")
    PPTX_AVAILABLE = False

# Try to import openpyxl for XLSX file processing
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    logger.warning("openpyxl not installed. Will not be able to read XLSX files.")
    OPENPYXL_AVAILABLE = False

# Try to import Playwright (but we'll have an alternative)
try:
    from playwright.sync_api import sync_playwright, TimeoutError as PlaywrightTimeoutError
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    logger.warning("Playwright not installed. Will use requests-based search instead.")
    PlaywrightTimeoutError = Exception  # Fallback definition if Playwright is not installed
    PLAYWRIGHT_AVAILABLE = False

# Import for the alternative search solution
try:
    import requests
    from bs4 import BeautifulSoup
    import html
    REQUESTS_AVAILABLE = True
except ImportError:
    logger.warning("Requests or BeautifulSoup not installed. Install with 'pip install requests beautifulsoup4'")
    REQUESTS_AVAILABLE = False

# List of common user agents for browser fingerprinting protection
USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:109.0) Gecko/20100101 Firefox/112.0",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.4 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36",
]

# Search engines to try
SEARCH_ENGINES = [
    {
        "name": "Google",
        "url": "https://www.google.com/search?q={query}",
        "results_selector": [".g", "[data-hveid]", ".MjjYud", "div[data-sokoban-container]"],
        "title_selector": ["h3", ".LC20lb", "[role='heading']"],
        "link_selector": ["a"],
        "snippet_selector": [".VwiC3b", ".s3v9rd", ".lEBKkf", ".s8bAkb"],
    },
    {
        "name": "Bing",
        "url": "https://www.bing.com/search?q={query}",
        "results_selector": [".b_algo"],
        "title_selector": ["h2"],
        "link_selector": ["a"],
        "snippet_selector": [".b_caption p"],
    },
    {
        "name": "DuckDuckGo",
        "url": "https://duckduckgo.com/?q={query}",
        "results_selector": [".result", ".result__body"],
        "title_selector": [".result__title", ".result__a"],
        "link_selector": [".result__a"],
        "snippet_selector": [".result__snippet"],
    }
]


class Tool(ABC):
    """Base class for all tools."""
    
    @property
    @abstractmethod
    def name(self) -> str:
        """The name of the tool, used for identifying it in requests."""
        pass
    
    @property
    @abstractmethod
    def description(self) -> str:
        """A description of what the tool does."""
        pass
    
    @abstractmethod
    def execute(self, **kwargs) -> Any:
        """Execute the tool with the given parameters."""
        pass


class FileCreationTool(Tool):
    """Tool for creating files in the OUTPUTS directory."""
    
    @property
    def name(self) -> str:
        return "create_file"
    
    @property
    def description(self) -> str:
        return "Create a file with the given content in the OUTPUTS directory."
    
    def execute(self, filename: str, content: str, file_type: Optional[str] = None) -> Dict[str, Any]:
        """Create a file with the given content in the OUTPUTS directory.
        
        Args:
            filename: The name of the file to create (without extension if file_type is provided)
            content: The content to write to the file
            file_type: Optional file type/extension (without the dot)
        
        Returns:
            A dictionary with the status of the operation and the path to the created file
        """
        logger.info(f"Creating file: {filename} with type {file_type or 'unspecified'}")
        
        try:
            # Create the OUTPUTS directory if it doesn't exist
            outputs_dir = "OUTPUTS"
            if not os.path.exists(outputs_dir):
                os.makedirs(outputs_dir)
                logger.info(f"Created OUTPUTS directory: {os.path.abspath(outputs_dir)}")
            
            # Clean the filename to ensure it's valid
            clean_filename = self._sanitize_filename(filename)
            
            # Add the file extension if provided
            if file_type:
                # Remove any leading dots from the file_type
                clean_file_type = file_type.lstrip('.')
                final_filename = f"{clean_filename}.{clean_file_type}"
            else:
                # If no file_type is provided, use the filename as is
                # (assuming it already has an extension)
                final_filename = clean_filename
                
                # Check if the filename has an extension, add .txt if not
                if '.' not in final_filename:
                    final_filename = f"{final_filename}.txt"
                    logger.info(f"No extension provided, adding .txt")
            
            # Create the full file path
            file_path = os.path.join(outputs_dir, final_filename)
            
            # Check if the file already exists
            if os.path.exists(file_path):
                # Add a timestamp to make the filename unique
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                name_part, ext_part = os.path.splitext(final_filename)
                final_filename = f"{name_part}_{timestamp}{ext_part}"
                file_path = os.path.join(outputs_dir, final_filename)
                logger.info(f"File already exists, creating with timestamp: {final_filename}")
            
            # Write the content to the file
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            logger.info(f"File created successfully: {file_path}")
            
            return {
                "status": "success",
                "message": f"File created successfully: {final_filename}",
                "file_path": file_path,
                "file_name": final_filename
            }
            
        except Exception as e:
            error_msg = f"Error creating file: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                "status": "error",
                "message": error_msg
            }
    
    def _sanitize_filename(self, filename: str) -> str:
        """Remove invalid characters from a filename."""
        # Replace invalid characters with underscores
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        
        # Trim whitespace
        filename = filename.strip()
        
        # Use a default name if the filename is empty after cleaning
        if not filename:
            filename = f"generated_file_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            logger.info(f"Empty filename after sanitization, using default: {filename}")
        
        return filename


class PlaywrightScreenshotTool:
    """Helper class to take screenshots with Playwright."""
    
    @staticmethod
    def take_screenshot(url: str, output_path: str, full_page: bool = True, timeout: int = 30000) -> bool:
        """Take a screenshot of a webpage using Playwright."""
        logger.info(f"Taking screenshot of {url} with Playwright")
        
        if not PLAYWRIGHT_AVAILABLE:
            logger.error("Playwright not available for taking screenshots")
            return False
        
        try:
            with sync_playwright() as p:
                # Launch a browser with a random user agent
                browser = p.chromium.launch(headless=True)
                context = browser.new_context(
                    viewport={"width": 1280, "height": 800},
                    user_agent=random.choice(USER_AGENTS)
                )
                
                # Create a new page and navigate to the URL
                page = context.new_page()
                page.set_default_timeout(timeout)
                
                try:
                    page.goto(url, wait_until="networkidle")
                except PlaywrightTimeoutError:
                    logger.warning(f"Timeout waiting for networkidle, continuing anyway")
                    # Try to wait a bit more
                    time.sleep(2)
                
                # Make sure the directory exists
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                
                # Take the screenshot
                page.screenshot(path=output_path, full_page=full_page)
                logger.info(f"Screenshot saved to {output_path}")
                
                # Close the browser
                browser.close()
                return True
                
        except Exception as e:
            logger.error(f"Error taking screenshot: {str(e)}")
            return False


class RequestsWebSearchTool(Tool):
    """Tool for searching the web using requests and BeautifulSoup - no browser automation needed."""
    
    @property
    def name(self) -> str:
        return "web_search"
    
    @property
    def description(self) -> str:
        return "Search the web for information on a given query."
    
    def execute(self, query: str) -> Dict[str, Any]:
        """Execute a web search using requests and BeautifulSoup."""
        logger.info(f"Starting requests-based web search for query: {query}")
        screenshots_dir = "search_screenshots"
        
        # Create screenshots directory if it doesn't exist
        if not os.path.exists(screenshots_dir):
            os.makedirs(screenshots_dir)
            
        try:
            # Check if requests and BeautifulSoup are available
            if not REQUESTS_AVAILABLE:
                logger.error("Requests or BeautifulSoup not installed")
                return {
                    "status": "error",
                    "message": "Requests or BeautifulSoup not installed. Install with 'pip install requests beautifulsoup4'"
                }
            
            all_results = []
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # Try search with Google
            google_results = self._search_google(query, screenshots_dir, timestamp)
            all_results.extend(google_results)
            
            # Try search with Bing if we don't have enough results
            if len(all_results) < 5:
                bing_results = self._search_bing(query, screenshots_dir, timestamp)
                all_results.extend(bing_results)
            
            # Try a news search if we're looking for recent information
            ddg_results = self._search_duckduckgo(query, screenshots_dir, timestamp)
            all_results.extend(ddg_results)
                
            # Filter out any duplicates
            unique_results = []
            seen_links = set()
            
            for result in all_results:
                if result["link"] not in seen_links:
                    seen_links.add(result["link"])
                    unique_results.append(result)
            
            logger.info(f"Web search completed with {len(unique_results)} unique results")
            
            # Create screenshots with Playwright if available
            if PLAYWRIGHT_AVAILABLE:
                search_urls = {
                    "google": f"https://www.google.com/search?q={urllib.parse.quote(query)}",
                    "bing": f"https://www.bing.com/search?q={urllib.parse.quote(query)}",
                    "ddg": f"https://duckduckgo.com/?q={urllib.parse.quote(query)}"
                }
                
                for engine, url in search_urls.items():
                    screenshot_path = os.path.join(screenshots_dir, f"{engine}_search_{timestamp}.png")
                    PlaywrightScreenshotTool.take_screenshot(url, screenshot_path)
            
            return {
                "status": "success",
                "results": unique_results,
                "query": query,
                "screenshots_dir": screenshots_dir
            }
            
        except Exception as e:
            error_msg = f"Error during web search: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg,
                "query": query
            }
    
    def _get_random_user_agent(self):
        """Get a random user agent to avoid detection."""
        return random.choice(USER_AGENTS)
    
    def _search_google(self, query: str, screenshots_dir: str, timestamp: str) -> List[Dict[str, Any]]:
        """Search Google using requests and BeautifulSoup."""
        results = []
        
        try:
            # Format the search URL
            search_url = f"https://www.google.com/search?q={urllib.parse.quote(query)}"
            logger.info(f"Searching Google for: {query}")
            
            # Make the request with a random user agent
            headers = {
                "User-Agent": self._get_random_user_agent(),
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                "Accept-Language": "en-US,en;q=0.9",
                "Referer": "https://www.google.com/"
            }
            
            response = requests.get(search_url, headers=headers, timeout=10)
            
            # Parse the HTML
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Find search result elements (this may need adjustment based on Google's layout)
            search_results = soup.select(".g")
            
            if not search_results:
                search_results = soup.select("[data-hveid]")
            
            if not search_results:
                search_results = soup.select(".MjjYud")
            
            logger.info(f"Found {len(search_results)} Google search results")
            
            # Process each result
            for result in search_results[:10]:  # Limit to 10 results
                # Find title
                title_elem = result.select_one("h3")
                
                # Find link
                link_elem = result.select_one("a")
                
                # Find snippet
                snippet_elem = result.select_one(".VwiC3b") or result.select_one(".s3v9rd")
                
                if title_elem and link_elem and "href" in link_elem.attrs:
                    title = title_elem.get_text().strip()
                    link = link_elem["href"]
                    
                    # Google prepends results with /url?q=
                    if link.startswith("/url?q="):
                        link = link.split("/url?q=")[1].split("&")[0]
                        link = urllib.parse.unquote(link)
                    
                    # Make sure it's a valid URL
                    if link.startswith("http"):
                        snippet = snippet_elem.get_text().strip() if snippet_elem else "No description available"
                        
                        results.append({
                            "title": title,
                            "link": link,
                            "snippet": snippet,
                            "source": "Google"
                        })
                        logger.info(f"Extracted Google result: {title}")
            
            return results
            
        except Exception as e:
            logger.error(f"Error during Google search: {str(e)}")
            return []
    
    def _search_bing(self, query: str, screenshots_dir: str, timestamp: str) -> List[Dict[str, Any]]:
        """Search Bing using requests and BeautifulSoup."""
        results = []
        
        try:
            # Format the search URL
            search_url = f"https://www.bing.com/search?q={urllib.parse.quote(query)}"
            logger.info(f"Searching Bing for: {query}")
            
            # Make the request with a random user agent
            headers = {
                "User-Agent": self._get_random_user_agent(),
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                "Accept-Language": "en-US,en;q=0.9"
            }
            
            response = requests.get(search_url, headers=headers, timeout=10)
            
            # Parse the HTML
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Find search result elements
            search_results = soup.select(".b_algo")
            logger.info(f"Found {len(search_results)} Bing search results")
            
            # Process each result
            for result in search_results[:5]:  # Limit to 5 results
                # Find title
                title_elem = result.select_one("h2")
                
                # Find link
                link_elem = result.select_one("a")
                
                # Find snippet
                snippet_elem = result.select_one(".b_caption p")
                
                if title_elem and link_elem and "href" in link_elem.attrs:
                    title = title_elem.get_text().strip()
                    link = link_elem["href"]
                    
                    # Make sure it's a valid URL
                    if link.startswith("http"):
                        snippet = snippet_elem.get_text().strip() if snippet_elem else "No description available"
                        
                        results.append({
                            "title": title,
                            "link": link,
                            "snippet": snippet,
                            "source": "Bing"
                        })
                        logger.info(f"Extracted Bing result: {title}")
            
            return results
            
        except Exception as e:
            logger.error(f"Error during Bing search: {str(e)}")
            return []
    
    def _search_duckduckgo(self, query: str, screenshots_dir: str, timestamp: str) -> List[Dict[str, Any]]:
        """Search DuckDuckGo using requests and BeautifulSoup."""
        results = []
        
        try:
            # DuckDuckGo's search API
            search_url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(query)}"
            logger.info(f"Searching DuckDuckGo for: {query}")
            
            # Make the request with a random user agent
            headers = {
                "User-Agent": self._get_random_user_agent(),
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                "Accept-Language": "en-US,en;q=0.9"
            }
            
            response = requests.get(search_url, headers=headers, timeout=10)
            
            # Parse the HTML
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Find search result elements
            search_results = soup.select(".result")
            logger.info(f"Found {len(search_results)} DuckDuckGo search results")
            
            # Process each result
            for result in search_results[:5]:  # Limit to 5 results
                # Find title and link
                title_elem = result.select_one(".result__title")
                link_elem = result.select_one(".result__url")
                
                # Find snippet
                snippet_elem = result.select_one(".result__snippet")
                
                if title_elem:
                    title_a = title_elem.select_one("a")
                    if title_a:
                        title = title_a.get_text().strip()
                        
                        # For DuckDuckGo, we need to extract the real URL
                        link = ""
                        if link_elem:
                            link = "https://" + link_elem.get_text().strip()
                        elif "href" in title_a.attrs:
                            href = title_a["href"]
                            if href.startswith("/"):
                                # Extract destination URL from DuckDuckGo redirect
                                redirect_match = re.search(r'uddg=([^&]+)', href)
                                if redirect_match:
                                    link = urllib.parse.unquote(redirect_match.group(1))
                        
                        if link and link.startswith("http"):
                            snippet = snippet_elem.get_text().strip() if snippet_elem else "No description available"
                            
                            results.append({
                                "title": title,
                                "link": link,
                                "snippet": snippet,
                                "source": "DuckDuckGo"
                            })
                            logger.info(f"Extracted DuckDuckGo result: {title}")
            
            return results
            
        except Exception as e:
            logger.error(f"Error during DuckDuckGo search: {str(e)}")
            return []
    
    def visit_and_summarize(self, url: str) -> Dict[str, Any]:
        """Visit a specific URL and extract content using requests and BeautifulSoup."""
        logger.info(f"Visiting and summarizing URL: {url}")
        screenshots_dir = "page_screenshots"
        
        if not os.path.exists(screenshots_dir):
            os.makedirs(screenshots_dir)
            
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        try:
            if not REQUESTS_AVAILABLE:
                return {"status": "error", "message": "Requests or BeautifulSoup not installed"}
            
            # Make the request with a random user agent
            headers = {
                "User-Agent": self._get_random_user_agent(),
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                "Accept-Language": "en-US,en;q=0.9"
            }
            
            response = requests.get(url, headers=headers, timeout=15)
            
            # Take a screenshot using Playwright if available
            screenshot_path = os.path.join(screenshots_dir, f"page_{timestamp}.png")
            screenshot_taken = False
            
            if PLAYWRIGHT_AVAILABLE:
                screenshot_taken = PlaywrightScreenshotTool.take_screenshot(url, screenshot_path)
                logger.info(f"Screenshot saved to {screenshot_path}" if screenshot_taken else "Failed to take screenshot")
            
            # Parse the HTML
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Extract page title
            title = soup.title.get_text().strip() if soup.title else "No title found"
            logger.info(f"Page title: {title}")
            
            # Try to get main content using various content selectors
            main_content = ""
            content_selectors = [
                "main", "article", "#content", ".content", 
                "[role='main']", ".main-content", ".post-content",
                "#main", ".article-content", ".entry-content", 
                "[itemprop='articleBody']", ".body", "#article-body"
            ]
            
            for selector in content_selectors:
                content_elem = soup.select_one(selector)
                if content_elem:
                    content_text = content_elem.get_text()
                    if len(content_text) > 100:  # Only use if substantial content
                        main_content = content_text
                        logger.info(f"Found content with selector: {selector}")
                        break
            
            # If no main content found, try to extract paragraphs
            if not main_content:
                logger.info("No main content found with selectors, extracting paragraphs")
                paragraphs = soup.select("p")
                paragraph_texts = []
                
                for p in paragraphs:
                    text = p.get_text().strip()
                    if len(text) > 40:  # Only include substantial paragraphs
                        paragraph_texts.append(text)
                
                main_content = "\n\n".join(paragraph_texts)
            
            # If still no content, get body text
            if not main_content:
                logger.info("Falling back to body text")
                body = soup.body
                if body:
                    main_content = body.get_text()
                else:
                    main_content = "Failed to extract content from page."
            
            # Try to extract metadata
            metadata = {}
            
            # Publication date
            date_selectors = [
                "time", "[itemprop='datePublished']", ".date", ".published", 
                "[datetime]", ".post-date", ".article-date"
            ]
            for selector in date_selectors:
                date_elem = soup.select_one(selector)
                if date_elem:
                    date_text = date_elem.get_text().strip()
                    if date_text:
                        metadata["publication_date"] = date_text
                        break
            
            # Author
            author_selectors = [
                "[itemprop='author']", ".author", ".byline", 
                "[rel='author']", ".article-author"
            ]
            for selector in author_selectors:
                author_elem = soup.select_one(selector)
                if author_elem:
                    author_text = author_elem.get_text().strip()
                    if author_text:
                        metadata["author"] = author_text
                        break
            
            # Clean and truncate content
            cleaned_content = self._clean_content(main_content)
            truncated_content = cleaned_content[:5000]  # Increased limit to 5000 chars
            
            result = {
                "status": "success",
                "title": title,
                "content": truncated_content,
                "metadata": metadata,
                "url": url
            }
            
            if screenshot_taken:
                result["screenshot"] = screenshot_path
                
            return result
                
        except Exception as e:
            error_msg = f"Error visiting URL: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg,
                "url": url
            }
    
    def _clean_content(self, content: str) -> str:
        """Clean extracted content by removing excess whitespace and common boilerplate."""
        if not content:
            return ""
        
        # Replace multiple newlines with just two
        content = "\n".join(line for line in content.splitlines() if line.strip())
        
        # Replace multiple spaces with single space
        content = " ".join(content.split())
        
        # Common phrases to remove (like cookie notices, navigation instructions)
        boilerplate = [
            "accept cookies", "cookie policy", "use cookies", 
            "privacy policy", "terms of service", "all rights reserved",
            "navigation menu", "skip to content", "search", "sign in",
            "subscribe to our newsletter", "subscribe now", "sign up",
            "we've updated our privacy policy"
        ]
        
        # Create a list of lines, filtering out boilerplate
        lines = content.splitlines()
        filtered_lines = []
        
        for line in lines:
            if line and not any(phrase in line.lower() for phrase in boilerplate):
                filtered_lines.append(line)
        
        return "\n".join(filtered_lines)


class ToolRegistry:
    """Registry for tools that can be used by the agent."""
    
    def __init__(self):
        self.tools = {}
    
    def register_tool(self, tool: Tool) -> None:
        """Register a tool with the registry."""
        self.tools[tool.name] = tool
    
    def get_tool(self, name: str) -> Tool:
        """Get a tool by name."""
        return self.tools.get(name)
    
    def has_tool(self, name: str) -> bool:
        """Check if a tool exists in the registry."""
        return name in self.tools
    
    def list_tools(self) -> Dict[str, str]:
        """List all available tools with their descriptions."""
        return {name: tool.description for name, tool in self.tools.items()}


class DocumentReaderTool(Tool):
    """Tool for reading the content of various document file formats."""
    
    @property
    def name(self) -> str:
        return "document_reader"
    
    @property
    def description(self) -> str:
        return "Read the content of various document files to analyze them and answer questions about them."
    
    def execute(self, file_path: str) -> Dict[str, Any]:
        """Read the content of a document file.
        
        Args:
            file_path: The path to the file to read
            
        Returns:
            A dictionary with the status of the operation and the content of the file
        """
        logger.info(f"Reading document: {file_path}")
        
        # Validate file exists
        if not os.path.exists(file_path):
            error_msg = f"File not found: {file_path}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
        
        # Get file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # Validate file type
        supported_extensions = [".pdf", ".txt", ".text", ".docx", ".pptx", ".xlsx", ".csv"]
        if ext not in supported_extensions:
            error_msg = f"Unsupported file type: {ext}. Supported formats: {', '.join(supported_extensions)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
            
        try:
            # Use the FileReaderTool to read the file content
            result = FileReaderTool.read_file(file_path)
            
            if result["status"] == "success":
                logger.info(f"Successfully read document: {file_path}")
                
                # Get file metadata
                file_stat = os.stat(file_path)
                file_size = file_stat.st_size
                file_modified = datetime.fromtimestamp(file_stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')
                
                # Add file info to result
                result["file_info"] = {
                    "path": file_path,
                    "name": os.path.basename(file_path),
                    "type": ext[1:].upper(),  # Remove dot and capitalize
                    "size": file_size,
                    "size_formatted": self._format_file_size(file_size),
                    "modified": file_modified
                }
                
            return result
                
        except Exception as e:
            error_msg = f"Error reading document: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                "status": "error",
                "message": error_msg
            }
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human-readable format."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024 or unit == 'GB':
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024


class FileReaderTool:
    """Static utility class for reading various document file formats."""
    
    @staticmethod
    def read_file(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read the content of a document file based on its extension.
        
        Supported formats:
        - PDF (.pdf) - Requires PyPDF2
        - Text (.txt, .text) - Plain text
        - Word (.docx) - Requires docx2txt
        - PowerPoint (.pptx) - Requires python-pptx
        - Excel (.xlsx) - Requires openpyxl
        - CSV (.csv) - Native Python
        """
        logger.info(f"Reading file: {file_path}")
        
        if not os.path.exists(file_path):
            return {"status": "error", "message": "File not found."}
        
        # Check file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # Process based on file type
        if ext == ".pdf":
            return FileReaderTool._read_pdf(file_path)
        elif ext in [".txt", ".text"]:
            return FileReaderTool._read_txt(file_path)
        elif ext == ".docx":
            return FileReaderTool._read_docx(file_path)
        elif ext == ".pptx":
            return FileReaderTool._read_pptx(file_path)
        elif ext == ".xlsx":
            return FileReaderTool._read_xlsx(file_path)
        elif ext == ".csv":
            return FileReaderTool._read_csv(file_path)
        else:
            return {"status": "error", "message": f"Unsupported file type: {ext}"}
    
    @staticmethod
    def _read_pdf(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from a PDF file."""
        if not PYPDF2_AVAILABLE:
            return {"status": "error", "message": "PyPDF2 not available. Install with 'pip install PyPDF2'"}
        
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                text = []
                
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text.append(f"--- Page {i+1} of {num_pages} ---\n{page_text}")
                
                return {
                    "status": "success",
                    "content": "\n\n".join(text),
                    "metadata": {
                        "num_pages": num_pages,
                        "has_text": len(text) > 0
                    }
                }
        
        except Exception as e:
            error_msg = f"Error reading PDF file: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
    
    @staticmethod
    def _read_txt(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from a TXT file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                # Count number of lines
                num_lines = content.count('\n') + 1
                
                return {
                    "status": "success",
                    "content": content,
                    "metadata": {
                        "num_lines": num_lines,
                        "size_chars": len(content)
                    }
                }
        
        except UnicodeDecodeError:
            # Try different encodings if UTF-8 fails
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    content = f.read()
                    num_lines = content.count('\n') + 1
                    
                    return {
                        "status": "success",
                        "content": content,
                        "metadata": {
                            "num_lines": num_lines,
                            "size_chars": len(content),
                            "encoding": "latin-1"
                        }
                    }
            except Exception as e:
                error_msg = f"Error reading TXT file with alternative encoding: {str(e)}"
                logger.error(error_msg)
                return {
                    "status": "error",
                    "message": error_msg
                }
        
        except Exception as e:
            error_msg = f"Error reading TXT file: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
    
    @staticmethod
    def _read_docx(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from a DOCX file."""
        if not DOCX2TXT_AVAILABLE:
            return {"status": "error", "message": "docx2txt not available. Install with 'pip install docx2txt'"}
        
        try:
            # Extract text from the docx
            text = docx2txt.process(file_path)
            
            # Count paragraphs (non-empty lines)
            paragraphs = [p for p in text.split('\n') if p.strip()]
            
            return {
                "status": "success",
                "content": text,
                "metadata": {
                    "num_paragraphs": len(paragraphs),
                    "size_chars": len(text)
                }
            }
        
        except Exception as e:
            error_msg = f"Error reading DOCX file: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
    
    @staticmethod
    def _read_pptx(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from a PPTX file."""
        if not PPTX_AVAILABLE:
            logger.error("python-pptx library not available")
            return {"status": "error", "message": "python-pptx not available. Install with 'pip install python-pptx'"}
        
        try:
            logger.info(f"Attempting to read PPTX file: {file_path}")
            
            # Add detailed logging for debugging purposes
            import os
            if not os.path.exists(file_path):
                logger.error(f"File not found: {file_path}")
                return {"status": "error", "message": f"File not found: {file_path}"}
            
            file_size = os.path.getsize(file_path)
            logger.info(f"PPTX file size: {file_size} bytes")
            
            # Try the default method to read the presentation
            try:
                from pptx import Presentation
                presentation = Presentation(file_path)
                
                # Log successful load
                logger.info(f"Successfully loaded presentation with {len(presentation.slides)} slides")
                
                slide_texts = []
                slide_count = len(presentation.slides)
                
                # Extract text from each slide
                for i, slide in enumerate(presentation.slides):
                    logger.info(f"Processing slide {i+1} of {slide_count}")
                    texts = []
                    
                    # Process all shapes in the slide
                    shape_count = len(slide.shapes)
                    logger.info(f"Slide {i+1} has {shape_count} shapes")
                    
                    for shape in slide.shapes:
                        # Log shape type
                        shape_type = type(shape).__name__
                        
                        # Check if it's a text-containing shape
                        if hasattr(shape, "text"):
                            if shape.text and len(shape.text.strip()) > 0:
                                texts.append(shape.text)
                                logger.info(f"Found text in {shape_type}: {shape.text[:50]}...")
                            else:
                                logger.info(f"Empty text in {shape_type}")
                        else:
                            logger.info(f"No text attribute in shape type: {shape_type}")
                    
                    # Also try to get text from text frames if available
                    if hasattr(slide, "shapes") and hasattr(slide.shapes, "text_frame"):
                        if slide.shapes.text_frame and hasattr(slide.shapes.text_frame, "text"):
                            texts.append(slide.shapes.text_frame.text)
                    
                    # For placeholders (like title, content)
                    for shape in slide.placeholders:
                        if hasattr(shape, "text") and shape.text:
                            texts.append(shape.text)
                    
                    # Assemble the slide text
                    if texts:
                        slide_text = "\n".join(texts)
                        slide_texts.append(f"--- Slide {i+1} of {slide_count} ---\n{slide_text}")
                    else:
                        slide_texts.append(f"--- Slide {i+1} of {slide_count} ---\n[No text content found]")
                
                return {
                    "status": "success",
                    "content": "\n\n".join(slide_texts),
                    "metadata": {
                        "num_slides": slide_count,
                        "slides_with_text": len([t for t in slide_texts if "[No text content found]" not in t])
                    }
                }
            
            except Exception as e:
                logger.error(f"Error using python-pptx Presentation: {str(e)}")
                logger.error("Falling back to alternative method...")
                
                # Alternative extraction method for problematic files
                try:
                    import zipfile
                    import re
                    from xml.etree import ElementTree
                    
                    # PPTX files are ZIP archives containing XML
                    slide_texts = []
                    with zipfile.ZipFile(file_path) as zf:
                        # Find all slide XML files
                        slide_files = [f for f in zf.namelist() if f.startswith('ppt/slides/slide')]
                        slide_files.sort()  # Ensure correct order
                        
                        logger.info(f"Found {len(slide_files)} slide files using zipfile method")
                        
                        # Process each slide
                        for i, slide_file in enumerate(slide_files):
                            try:
                                with zf.open(slide_file) as f:
                                    slide_xml = f.read()
                                    
                                # Parse the XML
                                root = ElementTree.fromstring(slide_xml)
                                
                                # Extract text using XPath-like search (simplified)
                                texts = []
                                # The XML namespace in PPTX files
                                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                
                                # Find all text elements
                                for text_elem in root.findall('.//a:t', ns):
                                    if text_elem.text and text_elem.text.strip():
                                        texts.append(text_elem.text)
                                
                                if texts:
                                    slide_text = "\n".join(texts)
                                    slide_texts.append(f"--- Slide {i+1} ---\n{slide_text}")
                                else:
                                    slide_texts.append(f"--- Slide {i+1} ---\n[No text content found]")
                                    
                            except Exception as xml_error:
                                logger.error(f"Error processing slide XML {slide_file}: {str(xml_error)}")
                                slide_texts.append(f"--- Slide {i+1} ---\n[Error extracting content]")
                    
                    return {
                        "status": "success",
                        "content": "\n\n".join(slide_texts),
                        "metadata": {
                            "num_slides": len(slide_files),
                            "slides_with_text": len([t for t in slide_texts if "[No text content found]" not in t]),
                            "extraction_method": "zipfile fallback"
                        }
                    }
                    
                except Exception as fallback_error:
                    logger.error(f"Fallback extraction method failed: {str(fallback_error)}")
                    raise  # Re-raise to be caught by the outer exception handler
        
        except Exception as e:
            error_msg = f"Error reading PPTX file: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                "status": "error",
                "message": error_msg
            }
    
    @staticmethod
    def _read_xlsx(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from an XLSX file."""
        if not OPENPYXL_AVAILABLE:
            return {"status": "error", "message": "openpyxl not available. Install with 'pip install openpyxl'"}
        
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheet_contents = []
            
            # Process each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                rows = []
                
                # Get data (limit to first 100 rows and 20 columns for performance)
                row_count = 0
                for row in sheet.iter_rows(max_row=100, max_col=20, values_only=True):
                    # Format the row and handle None values
                    formatted_row = [str(cell) if cell is not None else "" for cell in row]
                    rows.append("\t".join(formatted_row))
                    row_count += 1
                
                if rows:
                    sheet_text = f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows)
                    sheet_contents.append(sheet_text)
            
            return {
                "status": "success",
                "content": "\n\n".join(sheet_contents),
                "metadata": {
                    "num_sheets": len(workbook.sheetnames),
                    "sheet_names": workbook.sheetnames
                }
            }
        
        except Exception as e:
            error_msg = f"Error reading XLSX file: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }
    
    @staticmethod
    def _read_csv(file_path: str) -> Dict[str, Union[str, List[str]]]:
        """Read content from a CSV file."""
        try:
            import csv
            
            rows = []
            total_rows = 0
            
            with open(file_path, 'r', newline='', encoding='utf-8') as csvfile:
                csv_reader = csv.reader(csvfile)
                header = next(csv_reader, None)
                
                if header:
                    rows.append("\t".join(header))
                
                # Read up to 100 rows to avoid memory issues with large files
                for i, row in enumerate(csv_reader):
                    if i >= 100:
                        break
                    rows.append("\t".join(row))
                    total_rows = i + 1
            
            return {
                "status": "success",
                "content": "\n".join(rows),
                "metadata": {
                    "total_rows": total_rows + 1,  # Add 1 for header
                    "has_header": header is not None,
                    "num_columns": len(header) if header else 0,
                    "preview_note": "Showing first 100 rows only" if total_rows >= 100 else "Showing all rows"
                }
            }
            
        except UnicodeDecodeError:
            # Try different encodings if UTF-8 fails
            try:
                rows = []
                total_rows = 0
                
                with open(file_path, 'r', newline='', encoding='latin-1') as csvfile:
                    csv_reader = csv.reader(csvfile)
                    header = next(csv_reader, None)
                    
                    if header:
                        rows.append("\t".join(header))
                    
                    for i, row in enumerate(csv_reader):
                        if i >= 100:
                            break
                        rows.append("\t".join(row))
                        total_rows = i + 1
                
                return {
                    "status": "success",
                    "content": "\n".join(rows),
                    "metadata": {
                        "total_rows": total_rows + 1,
                        "has_header": header is not None,
                        "num_columns": len(header) if header else 0,
                        "encoding": "latin-1",
                        "preview_note": "Showing first 100 rows only" if total_rows >= 100 else "Showing all rows"
                    }
                }
            except Exception as e:
                error_msg = f"Error reading CSV file with alternative encoding: {str(e)}"
                logger.error(error_msg)
                return {
                    "status": "error",
                    "message": error_msg
                }
        
        except Exception as e:
            error_msg = f"Error reading CSV file: {str(e)}"
            logger.error(error_msg)
            return {
                "status": "error",
                "message": error_msg
            }